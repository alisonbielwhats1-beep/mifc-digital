from __future__ import annotations

import argparse
import json
import re
from pathlib import Path
from typing import Any


def normalize(text: str) -> str:
    return re.sub(r"\s+", " ", text or "").strip()


def extract_docx(path: Path) -> dict[str, Any]:
    from docx import Document

    document = Document(path)
    paragraphs = []
    for index, paragraph in enumerate(document.paragraphs, start=1):
        text = normalize(paragraph.text)
        if text:
            paragraphs.append(
                {
                    "location": f"paragraph:{index}",
                    "style": paragraph.style.name if paragraph.style else None,
                    "text": text,
                }
            )
    table_cells = []
    for table_index, table in enumerate(document.tables, start=1):
        for row_index, row in enumerate(table.rows, start=1):
            values = [normalize(cell.text) for cell in row.cells]
            if any(values):
                table_cells.append(
                    {
                        "location": f"table:{table_index}:row:{row_index}",
                        "text": " | ".join(values),
                    }
                )
    headings = [
        item for item in paragraphs
        if item["style"] and item["style"].lower().startswith(("heading", "título"))
    ]
    return {
        "kind": "docx",
        "path": str(path),
        "paragraph_count": len(document.paragraphs),
        "table_count": len(document.tables),
        "nonempty_paragraph_count": len(paragraphs),
        "table_row_count": len(table_cells),
        "headings": headings,
        "items": paragraphs + table_cells,
    }


def extract_pptx(path: Path) -> dict[str, Any]:
    from pptx import Presentation

    presentation = Presentation(path)
    slides = []
    items = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = normalize(shape.text)
                if text:
                    texts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    value = " | ".join(normalize(cell.text) for cell in row.cells)
                    if value.strip(" |"): 
                        texts.append(value)
        unique_texts = list(dict.fromkeys(texts))
        title = None
        if slide.shapes.title is not None:
            title = normalize(slide.shapes.title.text)
        combined = " || ".join(unique_texts)
        slides.append(
            {
                "slide": slide_index,
                "title": title,
                "text_count": len(unique_texts),
                "preview": combined[:500],
            }
        )
        if combined:
            items.append({"location": f"slide:{slide_index}", "text": combined})
    return {
        "kind": "pptx",
        "path": str(path),
        "slide_count": len(presentation.slides),
        "slides": slides,
        "items": items,
    }


def extract_pdf(path: Path) -> dict[str, Any]:
    from pypdf import PdfReader

    reader = PdfReader(path)
    pages = []
    items = []
    for page_index, page in enumerate(reader.pages, start=1):
        text = normalize(page.extract_text() or "")
        pages.append({"page": page_index, "preview": text[:500]})
        if text:
            items.append({"location": f"page:{page_index}", "text": text})
    return {
        "kind": "pdf",
        "path": str(path),
        "page_count": len(reader.pages),
        "pages": pages,
        "items": items,
    }


def extract(path: Path) -> dict[str, Any]:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        return extract_docx(path)
    if suffix == ".pptx":
        return extract_pptx(path)
    if suffix == ".pdf":
        return extract_pdf(path)
    raise ValueError(f"Unsupported file: {path}")


def search_items(result: dict[str, Any], terms: list[str]) -> list[dict[str, Any]]:
    patterns = [re.compile(term, re.IGNORECASE) for term in terms]
    matches = []
    for item in result["items"]:
        matched_terms = [term for term, pattern in zip(terms, patterns) if pattern.search(item["text"])]
        if matched_terms:
            matches.append({**item, "matched_terms": matched_terms})
    return matches


def summary(result: dict[str, Any]) -> dict[str, Any]:
    if result["kind"] == "docx":
        return {key: result[key] for key in ("kind", "path", "paragraph_count", "table_count", "nonempty_paragraph_count", "table_row_count", "headings")}
    if result["kind"] == "pptx":
        return {key: result[key] for key in ("kind", "path", "slide_count", "slides")}
    return {key: result[key] for key in ("kind", "path", "page_count", "pages")}


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("files", nargs="+", type=Path)
    parser.add_argument("--search", nargs="*")
    parser.add_argument("--full", action="store_true")
    args = parser.parse_args()

    outputs = []
    for path in args.files:
        result = extract(path)
        if args.search:
            outputs.append(
                {
                    "kind": result["kind"],
                    "path": result["path"],
                    "matches": search_items(result, args.search),
                }
            )
        elif args.full:
            outputs.append(result)
        else:
            outputs.append(summary(result))
    print(json.dumps(outputs, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
